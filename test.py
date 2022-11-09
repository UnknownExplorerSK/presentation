import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import openai
import time
import os
import requests
import shutil
import urllib

pr1 = Presentation()
#number of slides

var = input("Please enter your api key: ")
openai.api_key = var

var = input("How many slides? ")
i = int(var)
x = 0
var2 = input("The title of the presentation? ")
var3 = input("Who is making the presentation? ")
i = i - 1

# prompt2 = str(var)

while i > 0:
  i -= 1
  slide = "slide"
  title = "title"
  slide_number = (slide + str(x))
  title_number = (title + str(x))
  #print(slide_number)

  while x == 0:
    i -= 1
    slide_number_register = pr1.slide_layouts[0]
    slide_number = pr1.slides.add_slide(slide_number_register)

    title_number = slide_number.shapes.title
    title_number.text = str(var2)

    subtitle1 = slide_number.placeholders[1]
    subtitle1.text = ("Made by " + str(var3))
    x = x+1

  while x==1:
    i -= 1
    slide_number_register = pr1.slide_layouts[1]
    slide_number = pr1.slides.add_slide(slide_number_register)

    title_number = slide_number.shapes.title
    title_number.text = ("What is "+ str(var2.lower()) + "?")

    x = x+1

  else:
    slide_number_register = pr1.slide_layouts[1]
    slide_number = pr1.slides.add_slide(slide_number_register)

    title_number = slide_number.shapes.title
    var = input('Title of slide number ' + str(x+1) + ':')
    title_number.text = var

    prompt1 = str(var)
    response = openai.Completion.create(
     engine="text-davinci-001", 
     prompt=prompt1, 
     max_tokens=12)
    
    subtitle1 = slide_number.placeholders[1]
    subtitle1.text = str(response.choices[0].text)
    time.sleep(2)

    response2 = openai.Image.create(
     prompt=prompt1,
     n=1,
     size="256x256")
    image_url = response2['data'][0]['url']
    print(image_url)
    time.sleep(5)

    img_data = requests.get(image_url).content
    with open('dall_e_generated_image.png', 'wb') as handler:
        handler.write(img_data)
    time.sleep(5)

    img1 = "dall_e_generated_image.png"

    from_left = Inches(4)
    from_top = Inches(4)
    add_picture = slide_number.shapes.add_picture(img1,from_left,from_top)
    time.sleep(2)

    x = x+1

  while i == 2:
    i -= 1
    slide_number_register = pr1.slide_layouts[1]
    slide_number = pr1.slides.add_slide(slide_number_register)

    title_number = slide_number.shapes.title
    title_number.text = (str(var2) + " in the future")

    x = x+1
    
  while i == 1:
    i -= 1
    slide_number_register = pr1.slide_layouts[1]
    slide_number = pr1.slides.add_slide(slide_number_register)

    title_number = slide_number.shapes.title
    title_number.text = "Used sources"

    x = x+1

  while i == 0:
    slide_number_register = pr1.slide_layouts[0]
    slide_number = pr1.slides.add_slide(slide_number_register)

    title_number = slide_number.shapes.title
    title_number.text = "Are there any questions?"

    subtitle1 = slide_number.placeholders[1]
    subtitle1.text = "Have a wonderful day!"
    x = x+1
    i -= 1
  
else:
  pr1.save("Testing.pptx")
  print("Presentation done!")
