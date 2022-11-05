import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

pr1 = Presentation()

slide1_register = pr1.slide_layouts[0]
slide1 = pr1.slides.add_slide(slide1_register)

title1 = slide1.shapes.title
title1.text = "Censorship"

subtitle1 = slide1.placeholders[1]
subtitle1.text = "Made by Sebastians"

slide2_register = pr1.slide_layouts[1]
slide2 = pr1.slides.add_slide(slide2_register)

title2 = slide2.shapes.title
title2.text = "Propoganda"

bullet_point_box = slide2.shapes

bullet_points_lvl1 = bullet_point_box.placeholders[1]
bullet_points_lvl1.text = "1st point"

bullet_points_lvl2 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl2.text = "2nd point"
bullet_points_lvl2.level = 1

bullet_points_lvl3 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl3.text = "3rd point"
bullet_points_lvl3.level = 2

bullet_points_lvl4 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl4.text = "4th point"
bullet_points_lvl4.level = 3

sliede3_register = pr1.slide_layouts[5]
slide3 = pr1.slides.add_slide(sliede3_register)

title3 = slide3.shapes.title
title3.text = "Pictures!!!"

image = Image.open('pexels-pixabay-414102.jpg')
new_image = image.resize((400,200))
new_image.save('newimage_500.jpg')

img1 = "newimage_500.jpg"

from_left = Inches(2)
from_top = Inches(2)
add_picture = slide3.shapes.add_picture(img1,from_left,from_top)


pr1.save("Testing.pptx")